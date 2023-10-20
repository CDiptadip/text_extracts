from flask import Flask, request, render_template

from werkzeug.utils import secure_filename

import os

from PyPDF2 import PdfReader

from docx import Document

from openpyxl import load_workbook, Workbook

from pptx import Presentation

 

app = Flask(__name__)

app.config['UPLOAD_FOLDER'] = 'uploads'

app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'xlsx', 'pptx'}

data_to_store = []  # List to store data from each file

 

def allowed_file(filename):

    return '.' in filename and filename.rsplit('.', 1)[1] in app.config['ALLOWED_EXTENSIONS']

 

def extract_text_from_pdf(filename):

    text = ""

    with open(filename, 'rb') as pdf_file:

        pdf_reader = PdfReader(pdf_file)

        for page in pdf_reader.pages:

            text += page.extract_text()

    return text

 

def extract_text_from_docx(filename):

    doc = Document(filename)

    text = []

    for para in doc.paragraphs:

        text.append(para.text)

    return '\n'.join(text)

 

def extract_data_from_excel(filename):

    data = []

    wb = load_workbook(filename)

    ws = wb.active

    for row in ws.iter_rows(values_only=True):

        data.append(row)

    return data

 

def extract_text_from_pptx(filename):

    presentation = Presentation(filename)

    text = []

 

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text.append(shape.text)

 

    return "\n".join(text)

 

def save_to_excel():

    if data_to_store:

        excel_file = os.path.join(app.config['UPLOAD_FOLDER'], 'extracted_data.xlsx')

        wb = Workbook()

        ws = wb.active

        ws.append(["File Name", "Extracted Text"])

       

        for data in data_to_store:

            ws.append(data)

 

        wb.save(excel_file)

 

@app.route('/', methods=['GET', 'POST'])

def index():

    if request.method == 'POST':

        if 'file' not in request.files:

            return "No file part"

        file = request.files['file']

        if file.filename == '':

            return "No selected file"

        if file and allowed_file(file.filename):

            filename = secure_filename(file.filename)

            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)

            file.save(file_path)

 

            text = ""

            data = None

            pptx_text = None

 

            if filename.endswith('.pdf'):

                text = extract_text_from_pdf(file_path)

            elif filename.endswith('.docx'):

                text = extract_text_from_docx(file_path)

            elif filename.endswith('.xlsx'):

                data = extract_data_from_excel(file_path)

            elif filename.endswith('.pptx'):

                pptx_text = extract_text_from_pptx(file_path)

 

            data_to_store.append([filename, text])

 

            save_to_excel()  # Save to Excel after processing each file

 

            return render_template('index.html', text=text, data=data, pptx_text=pptx_text)

    return render_template('index.html', text=None, data=None, pptx_text=None)

 

if __name__ == '__main__':

    app.run(debug=True)